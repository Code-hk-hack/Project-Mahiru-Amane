import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

def add_slide(title_text, subtitle_text, bullets):
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    # Customize title font
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Arial'
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 153) # Tech blue
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    # Add subtitle
    p = tf.paragraphs[0]
    p.text = subtitle_text
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(100, 100, 100)
    
    # Add bullets
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(16)
        p.font.name = 'Arial'

# Slide 1
add_slide(
    "The Problem",
    "Social Anxiety, Poor Interview Skills, and No Affordable Real-Time Coaching",
    [
        "40% of Gen Zs feel stressed/anxious all or most of the time.",
        "74% of young professionals have needed time off due to stress, but fear speaking up.",
        "Career coaching costs $150-$300/hour—unaffordable for early-career professionals.",
        "Traditional coaching lacks instant, real-time feedback on speech patterns and passiveness.",
        "Result: Lost job opportunities, low confidence, and prolonged workplace anxiety."
    ]
)

# Slide 2
add_slide(
    "The Solution – Meet Mahiru/Amane",
    "Your 24/7 AI Coach for Social Confidence & Interview Mastery",
    [
        "Real-time voice-to-voice AI coaching simulator—speak naturally, get instant feedback.",
        "AI evaluates communication skills: detects passiveness, hesitation, filler words, and tone.",
        "Provides instant structural feedback on speech patterns and conversational flow.",
        "Completely judgment-free zone—practice as much as you want, anytime.",
        "Builds confidence through repeated, safe practice with actionable feedback."
    ]
)

# Slide 3
add_slide(
    "Technology Stack",
    "Built with Modern, Scalable Tech for Real-Time AI Coaching",
    [
        "Frontend: Next.js, React, Tailwind CSS, Aceternity UI",
        "Backend: Python, FastAPI, WebSockets (for real-time voice streaming)",
        "Database: Supabase (PostgreSQL) – session history, user progress tracking",
        "AI Engine: Groq (Llama-3 70B) / Google Gemini – core conversational intelligence",
        "Voice: Gnani.ai – ultra-fast Speech-to-Text (STT) & Text-to-Speech (TTS)",
        "Extensibility: MCP (Model Context Protocol) – dynamic tool execution"
    ]
)

# Slide 4
add_slide(
    "Strategic Partner Integrations",
    "Powered by Industry-Leading Hackathon Partners",
    [
        "Gnani.ai: Real-time STT/TTS enabling ultra-fast, natural voice conversations (<500ms latency).",
        "Groq / Google Gemini: Core LLM intelligence for nuanced coaching and feedback.",
        "Supabase: Secure database management, session history, and user analytics.",
        "Slashy (via MCP): AI Agent dynamically executes external tools and tasks during coaching."
    ]
)

# Slide 5
add_slide(
    "Market Opportunity",
    "A $236B+ EdTech Market Meets Demand for Social Confidence Training",
    [
        "Global EdTech Market: Growing to $236.3B (2026) at 18.3% CAGR.",
        "AI Career Coach Market: Explosive growth reaching $14.82B by 2030.",
        "Gen Z Mental Health Crisis: Mental health is the #2 societal concern for Gen Z.",
        "Target Audience: 2.3B Gen Z globally, entering the workforce.",
        "Market Gap: No affordable, real-time, AI-powered simulator exists.",
        "Bottom Line: Positioned perfectly at the intersection of EdTech, AI Coaching, and Workplace Mental Health."
    ]
)

# Save
prs.save("Project_Mahiru_Pitch_Deck.pptx")
print("Successfully generated Project_Mahiru_Pitch_Deck.pptx")
